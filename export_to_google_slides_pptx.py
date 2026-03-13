from __future__ import annotations

import http.server
import shutil
import socketserver
import subprocess
import threading
from functools import partial
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Emu


ROOT = Path(__file__).resolve().parent
HTML_PATH = ROOT / "magic_artist_deck_v3.html"
OUTPUT_DIR = ROOT / "artifacts"
RENDER_DIR = OUTPUT_DIR / "rendered_slides"
PPTX_PATH = OUTPUT_DIR / "magic_artist_deck_google_slides_import.pptx"
BROWSER_PROFILE_DIR = OUTPUT_DIR / "browser_profile"
VIEWPORT = (1484, 812)
EMU_PER_PX = 9525


def find_browser() -> Path:
  candidates = [
    Path(r"C:\Program Files (x86)\Microsoft\Edge\Application\msedge.exe"),
    Path(r"C:\Program Files\Google\Chrome\Application\chrome.exe"),
  ]
  for candidate in candidates:
    if candidate.exists():
      return candidate

  for name in ("msedge", "chrome"):
    resolved = shutil.which(name)
    if resolved:
      return Path(resolved)

  raise RuntimeError("Could not find Edge or Chrome for headless slide rendering.")


def slide_count() -> int:
  soup = BeautifulSoup(HTML_PATH.read_text(encoding="utf-8"), "lxml")
  return len(soup.select(".slide"))


class QuietHandler(http.server.SimpleHTTPRequestHandler):
  def log_message(self, format: str, *args) -> None:
    return


def start_server() -> tuple[socketserver.TCPServer, int]:
  handler = partial(QuietHandler, directory=str(ROOT))
  server = socketserver.TCPServer(("127.0.0.1", 0), handler)
  port = int(server.server_address[1])
  thread = threading.Thread(target=server.serve_forever, daemon=True)
  thread.start()
  return server, port


def render_slide(browser: Path, port: int, number: int, output_path: Path) -> None:
  width, height = VIEWPORT
  url = f"http://127.0.0.1:{port}/{HTML_PATH.name}?slide={number}&export=1"
  BROWSER_PROFILE_DIR.mkdir(parents=True, exist_ok=True)
  base_args = [
    str(browser),
    "--disable-gpu",
    "--disable-crash-reporter",
    "--disable-crashpad-for-testing",
    "--hide-scrollbars",
    "--no-sandbox",
    "--no-default-browser-check",
    "--no-first-run",
    "--force-device-scale-factor=1",
    f"--user-data-dir={BROWSER_PROFILE_DIR}",
    f"--window-size={width},{height}",
    f"--screenshot={output_path}",
    "--run-all-compositor-stages-before-draw",
    "--virtual-time-budget=1500",
    url,
  ]

  commands = [
    [base_args[0], "--headless=new", *base_args[1:]],
    [base_args[0], "--headless", *base_args[1:]],
  ]

  last_error: RuntimeError | None = None
  for command in commands:
    try:
      subprocess.run(command, check=True, timeout=30, stdout=subprocess.PIPE, stderr=subprocess.PIPE)
      return
    except subprocess.CalledProcessError as exc:
      last_error = RuntimeError(exc.stderr.decode("utf-8", errors="ignore") or exc.stdout.decode("utf-8", errors="ignore"))
    except subprocess.TimeoutExpired as exc:
      last_error = RuntimeError(f"Rendering timed out for slide {number}: {exc}")

  raise last_error or RuntimeError(f"Failed to render slide {number}.")


def build_pptx(image_paths: list[Path]) -> None:
  width, height = VIEWPORT
  presentation = Presentation()
  presentation.slide_width = Emu(width * EMU_PER_PX)
  presentation.slide_height = Emu(height * EMU_PER_PX)
  blank = presentation.slide_layouts[6]

  for image_path in image_paths:
    slide = presentation.slides.add_slide(blank)
    slide.shapes.add_picture(
      str(image_path),
      0,
      0,
      width=presentation.slide_width,
      height=presentation.slide_height,
    )

  OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
  presentation.save(str(PPTX_PATH))


def main() -> None:
  if not HTML_PATH.exists():
    raise FileNotFoundError(f"Deck HTML not found: {HTML_PATH}")

  browser = find_browser()
  total = slide_count()
  OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
  RENDER_DIR.mkdir(parents=True, exist_ok=True)

  server, port = start_server()
  try:
    rendered: list[Path] = []
    for number in range(1, total + 1):
      image_path = RENDER_DIR / f"slide-{number:02d}.png"
      print(f"Rendering slide {number}/{total} -> {image_path.name}")
      render_slide(browser, port, number, image_path)
      rendered.append(image_path)

    print(f"Building PPTX -> {PPTX_PATH.name}")
    build_pptx(rendered)
  finally:
    server.shutdown()
    server.server_close()

  print(f"Done: {PPTX_PATH}")


if __name__ == "__main__":
  main()
